#!/usr/bin/env python3
"""
Deep layout extraction — captures shape types, positions, fills,
and text container metadata so we know WHERE copy sits and in
WHAT component (pill, card, stat, bullet, etc).
"""

import json
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.util import Emu
from pptx.dml.color import RGBColor


EMU_PER_IN = 914400


def emu_to_pct(emu, total):
    if emu is None or total is None or total == 0:
        return None
    return round(float(emu) / float(total) * 100, 2)


def get_fill(shape):
    """Best-effort fill color extraction (solid fills only)."""
    try:
        if not shape.has_text_frame and shape.shape_type not in (
            MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER
        ):
            pass
        fill = shape.fill
        if fill.type == 1:  # MSO_FILL.SOLID
            try:
                return "#" + str(fill.fore_color.rgb)
            except Exception:
                return None
    except Exception:
        return None
    return None


def get_autoshape_type(shape):
    """Return a friendly name for auto-shapes (rect, roundRect, ellipse, etc)."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            ast = shape.auto_shape_type
            return str(ast).split(' (')[0] if ast else "AUTO"
    except Exception:
        pass
    return None


def classify_shape(shape):
    """Classify a shape into a semantic component type."""
    st = shape.shape_type

    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "placeholder"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            ast = shape.auto_shape_type
            ast_name = str(ast)
            if "ROUNDED_RECTANGLE" in ast_name or "ROUND" in ast_name:
                return "pill" if shape.has_text_frame else "round_rect"
            if "RECTANGLE" in ast_name:
                return "rect"
            if "OVAL" in ast_name or "ELLIPSE" in ast_name:
                return "oval"
            if "LINE" in ast_name:
                return "line"
            if "ARROW" in ast_name:
                return "arrow"
            return "autoshape"
        except Exception:
            return "autoshape"
    if st == MSO_SHAPE_TYPE.LINE:
        return "line"
    if st == MSO_SHAPE_TYPE.FREEFORM:
        return "freeform"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    return f"other_{st}"


def get_text_runs(shape):
    """Extract text with per-run formatting (font, size, color, bold)."""
    if not shape.has_text_frame:
        return []
    runs_out = []
    for para in shape.text_frame.paragraphs:
        para_runs = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            info = {"text": text}
            try:
                if run.font.size:
                    info["size_pt"] = run.font.size.pt
            except Exception:
                pass
            try:
                if run.font.name:
                    info["font"] = run.font.name
            except Exception:
                pass
            try:
                if run.font.bold:
                    info["bold"] = True
            except Exception:
                pass
            try:
                if run.font.italic:
                    info["italic"] = True
            except Exception:
                pass
            try:
                color = run.font.color
                if color.type is not None:
                    info["color"] = "#" + str(color.rgb)
            except Exception:
                pass
            para_runs.append(info)
        if para_runs:
            alignment = None
            try:
                if para.alignment is not None:
                    alignment = str(para.alignment)
            except Exception:
                pass
            runs_out.append({"runs": para_runs, "alignment": alignment})
    return runs_out


def walk(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(shape.shapes)
        else:
            yield shape


def extract(file_path, output_path):
    prs = Presentation(file_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    slides_out = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        shapes_out = []
        for shape in walk(slide.shapes):
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
            except Exception:
                left = top = width = height = None

            kind = classify_shape(shape)
            fill = get_fill(shape)
            autoshape = get_autoshape_type(shape)

            shape_data = {
                "kind": kind,
                "x_pct": emu_to_pct(left, slide_w),
                "y_pct": emu_to_pct(top, slide_h),
                "w_pct": emu_to_pct(width, slide_w),
                "h_pct": emu_to_pct(height, slide_h),
                "fill": fill,
                "autoshape": autoshape,
                "name": shape.name if hasattr(shape, "name") else None,
            }

            if shape.has_text_frame:
                shape_data["text"] = shape.text_frame.text
                shape_data["paragraphs"] = get_text_runs(shape)

            if kind == "image":
                shape_data["image_ref"] = None  # we don't re-extract here

            shapes_out.append(shape_data)

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
        except Exception:
            pass

        slides_out.append({
            "number": slide_num,
            "slide_w_emu": slide_w,
            "slide_h_emu": slide_h,
            "shapes": shapes_out,
            "notes": notes,
        })

    with open(output_path, "w") as f:
        json.dump(slides_out, f, indent=2, default=str)

    return slides_out


if __name__ == "__main__":
    slides = extract(sys.argv[1], sys.argv[2])
    print(f"Extracted {len(slides)} slides")
    for s in slides[:3]:
        print(f"  Slide {s['number']}: {len(s['shapes'])} shapes")
        kinds = {}
        for sh in s['shapes']:
            kinds[sh['kind']] = kinds.get(sh['kind'], 0) + 1
        print(f"    Kinds: {kinds}")
