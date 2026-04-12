#!/usr/bin/env python3
"""
Robust PPTX extractor. Handles placeholder/linked pictures, group shapes,
and picture-style shape fills that the stock script crashes on.
"""

import json
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def walk_shapes(shapes):
    """Recursively yield every shape, flattening group shapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)
        else:
            yield shape


def extract_pptx(file_path, output_dir="."):
    prs = Presentation(file_path)
    slides_data = []

    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    skipped_images = 0
    total_images = 0

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "number": slide_num + 1,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
        }

        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None

        for shape in walk_shapes(slide.shapes):
            # Text
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if title_shape is not None and shape == title_shape and not slide_data["title"]:
                    slide_data["title"] = text
                else:
                    slide_data["content"].append({"type": "text", "content": text})

            # Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_images += 1
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                except Exception as e:
                    skipped_images += 1
                    continue

                image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                image_path = os.path.join(assets_dir, image_name)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                try:
                    w = shape.width
                    h = shape.height
                    l = shape.left
                    t = shape.top
                except Exception:
                    w = h = l = t = None

                slide_data["images"].append(
                    {
                        "path": f"assets/{image_name}",
                        "width": w,
                        "height": h,
                        "left": l,
                        "top": t,
                    }
                )

        # Speaker notes
        try:
            if slide.has_notes_slide:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text
        except Exception:
            pass

        slides_data.append(slide_data)

    output_path = os.path.join(output_dir, "extracted-slides.json")
    with open(output_path, "w") as f:
        json.dump(slides_data, f, indent=2)

    return slides_data, total_images, skipped_images


if __name__ == "__main__":
    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    slides, total, skipped = extract_pptx(input_file, output_dir)

    print(f"\n=== EXTRACTION COMPLETE ===")
    print(f"Slides:          {len(slides)}")
    print(f"Total pictures:  {total}")
    print(f"Skipped (no embedded bytes): {skipped}")
    print(f"Saved images:    {total - skipped}")
    print()
    for s in slides:
        img_count = len(s["images"])
        title = s["title"] or "(no title)"
        if len(title) > 60:
            title = title[:57] + "..."
        print(f"  Slide {s['number']:3d}: {title:60s} — {img_count} img")
